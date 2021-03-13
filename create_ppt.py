from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import date
import os
import json

with open("data.json") as json_data:
    data_loaded = json.load(json_data)
    # print(data_loaded)
    include_name = data_loaded["include_name"]
    ppt_out_path = data_loaded["out_path"]
    material_list = data_loaded["material_list"]
    material_list = list(dict.fromkeys(material_list))


# Functions go here
def create_ppt(ppt_out_path):
    # use inbuilt python ppt Template
    prs = Presentation()

    # Create a title slide first
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Material Report"
    subtitle.text = "{} \n Generated on {:%m-%d-%Y}".format(include_name, date.today())

    for mat in material_list:
        material_name = mat
        image_name = material_name + ".png"

        # add a new slide
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # add picture
        img_path = os.path.join(ppt_out_path, image_name)
        left = Inches(3)
        top = Inches(2)
        height = Inches(4)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)

        # add text
        left = top = width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        # add material name
        p = tf.add_paragraph()
        p.text = "Material: {}".format(material_name)
        p.font.size = Pt(12)

    # Write PPT Out
    prs.save("{}/Mat_Report.pptx".format(ppt_out_path))


create_ppt(ppt_out_path=ppt_out_path)



